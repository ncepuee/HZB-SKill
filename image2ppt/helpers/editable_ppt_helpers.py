"""
Reusable helpers for editable PowerPoint reconstruction.
Requires: python-pptx
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def pixel_to_ppt(x, y, w, h, img_w, img_h, slide_w_in, slide_h_in):
    """Map image pixel coordinates to PowerPoint inches."""
    return (
        x / img_w * slide_w_in,
        y / img_h * slide_h_in,
        w / img_w * slide_w_in,
        h / img_h * slide_h_in,
    )


def add_text(slide, x, y, w, h, text, *,
             size=12, bold=False, color="#0E2240",
             font="Microsoft YaHei",
             align=PP_ALIGN.LEFT,
             valign=MSO_VERTICAL_ANCHOR.MIDDLE,
             margin=0.02):
    tb = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return tb


def add_box(slide, x, y, w, h, *,
            line_color="#B8C3D1",
            fill_color="#FFFFFF",
            radius=True,
            line_width=1.0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill_color)
    shp.line.color.rgb = rgb(line_color)
    shp.line.width = Pt(line_width)
    return shp


def add_circle(slide, cx, cy, r, *,
               fill_color="#1148E8",
               line_color=None,
               line_width=0.5):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx-r), Inches(cy-r), Inches(2*r), Inches(2*r)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill_color)
    if line_color:
        shp.line.color.rgb = rgb(line_color)
        shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    return shp


def add_line(slide, x1, y1, x2, y2, *,
             color="#0E2240",
             width=1.2):
    ln = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(width)
    return ln


def validate_totals(rows, keys=("r1", "r2")):
    """Return round totals and grand total before rendering."""
    totals = {k: sum(row.get(k, 0) for row in rows) for k in keys}
    totals["grand_total"] = sum(totals.values())
    return totals
