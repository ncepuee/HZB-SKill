#!/usr/bin/env python3
"""Structural and editability audit for reconstructed PowerPoint files."""

from __future__ import annotations

import argparse
import json
import unicodedata
import zipfile
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def expected_text_from_scene(path: str) -> list[list[str]]:
    if not path:
        return []
    scene = json.loads(Path(path).read_text(encoding="utf-8"))
    if "slides" in scene:
        return [slide.get("expected_text", []) for slide in scene["slides"]]
    return [scene.get("expected_text", [])]


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--scene", default="")
    ap.add_argument("--json-out", default="")
    ap.add_argument("--large-picture-threshold", type=float, default=0.85)
    args = ap.parse_args()

    pptx = Path(args.pptx).resolve()
    errors: list[str] = []
    warnings: list[str] = []
    media = []
    try:
        with zipfile.ZipFile(pptx) as zf:
            bad = zf.testzip()
            if bad:
                errors.append(f"corrupt package member: {bad}")
            for info in zf.infolist():
                if info.filename.startswith("ppt/media/"):
                    media.append({"name": info.filename, "size": info.file_size})
                    if info.file_size == 0:
                        errors.append(f"zero-byte media: {info.filename}")
    except Exception as exc:
        errors.append(f"cannot inspect package: {exc}")

    prs = Presentation(pptx)
    sw, sh = prs.slide_width, prs.slide_height
    expected_by_slide = expected_text_from_scene(args.scene)
    slides = []
    for si, slide in enumerate(prs.slides, 1):
        types = Counter()
        native_text_parts = []
        slide_errors = []
        slide_warnings = []
        max_picture_coverage = 0.0
        for idx, shape in enumerate(slide.shapes, 1):
            types[str(shape.shape_type)] += 1
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > sw or shape.top + shape.height > sh:
                slide_errors.append(f"shape {idx} out of bounds: {shape.name}")
            if getattr(shape, "has_text_frame", False):
                text = unicodedata.normalize("NFC", shape.text or "")
                if text.strip():
                    native_text_parts.append(text)
                elif shape.shape_type in {MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER}:
                    slide_warnings.append(f"empty explicit text shape {idx}: {shape.name}")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                coverage = (shape.width * shape.height) / (sw * sh)
                max_picture_coverage = max(max_picture_coverage, coverage)
                if coverage > args.large_picture_threshold:
                    slide_warnings.append(f"large picture {idx}: {coverage:.1%} coverage")

        corpus = "\n".join(native_text_parts)
        expected = expected_by_slide[si - 1] if si <= len(expected_by_slide) else []
        for item in expected:
            normalized = unicodedata.normalize("NFC", str(item))
            if normalized not in corpus:
                slide_errors.append(f"expected native text missing: {item}")
        if len(slide.shapes) == 1 and max_picture_coverage > args.large_picture_threshold:
            slide_errors.append("slide is effectively a single large picture")

        errors.extend([f"slide {si}: {x}" for x in slide_errors])
        warnings.extend([f"slide {si}: {x}" for x in slide_warnings])
        slides.append({
            "slide": si,
            "shape_count": len(slide.shapes),
            "shape_types": dict(types),
            "native_text_shapes": len(native_text_parts),
            "native_text_characters": len(corpus),
            "picture_count": sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE),
            "max_picture_coverage": round(max_picture_coverage, 6),
            "table_count": sum(1 for s in slide.shapes if getattr(s, "has_table", False)),
            "chart_count": sum(1 for s in slide.shapes if getattr(s, "has_chart", False)),
            "errors": slide_errors,
            "warnings": slide_warnings,
        })

    report = {
        "file": str(pptx),
        "slides": slides,
        "media": media,
        "errors": errors,
        "warnings": warnings,
        "passed": not errors,
    }
    if args.json_out:
        out = Path(args.json_out)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 1 if errors else 0


if __name__ == "__main__":
    raise SystemExit(main())

