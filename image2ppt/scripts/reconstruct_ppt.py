#!/usr/bin/env python3
"""
reconstruct_ppt.py - Helper toolkit for reconstructing image layouts into native, editable PowerPoint (.pptx).
Supports coordinate normalization, color tokens, cards, shapes, native tables, and text styling.
"""
import sys
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip('#')
    if len(hex_str) == 6:
        r = int(hex_str[0:2], 16)
        g = int(hex_str[2:4], 16)
        b = int(hex_str[4:6], 16)
        return RGBColor(r, g, b)
    return RGBColor(0, 0, 0)

class PptBuilder:
    def __init__(self, width_inches: float = 13.333, height_inches: float = 7.5, ref_w: float = 1920.0, ref_h: float = 1080.0):
        """
        Initialize 16:9 (default 13.333 x 7.5 inches) presentation with reference pixel mapping.
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(width_inches)
        self.prs.slide_height = Inches(height_inches)
        self.slide_width = width_inches
        self.slide_height = height_inches
        self.ref_w = float(ref_w)
        self.ref_h = float(ref_h)
        
        blank_slide_layout = self.prs.slide_layouts[6]
        self.slide = self.prs.slides.add_slide(blank_slide_layout)

    def map_x(self, px_x: float) -> Inches:
        return Inches(px_x / self.ref_w * self.slide_width)

    def map_y(self, px_y: float) -> Inches:
        return Inches(px_y / self.ref_h * self.slide_height)

    def map_w(self, px_w: float) -> Inches:
        return Inches(px_w / self.ref_w * self.slide_width)

    def map_h(self, px_h: float) -> Inches:
        return Inches(px_h / self.ref_h * self.slide_height)

    def add_card(self, x, y, w, h, bg_color="#FFFFFF", border_color="#E0E0E0", border_width=1.0, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
        """Add a background card/container with border and fill."""
        left, top, width, height = self.map_x(x), self.map_y(y), self.map_w(w), self.map_h(h)
        shape = self.slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(bg_color)
        if border_color:
            shape.line.color.rgb = hex_to_rgb(border_color)
            shape.line.width = Pt(border_width)
        else:
            shape.line.fill.background()
        return shape

    def add_text(self, x, y, w, h, text: str, font_size: float = 14, font_name: str = "Microsoft YaHei", 
                 color: str = "#333333", bold: bool = False, align=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.TOP, word_wrap=True):
        """Add an editable text box."""
        left, top, width, height = self.map_x(x), self.map_y(y), self.map_w(w), self.map_h(h)
        txBox = self.slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        tf.vertical_anchor = vertical_anchor
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        lines = text.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.alignment = align
            p.font.name = font_name
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = hex_to_rgb(color)
        return txBox

    def add_table(self, x, y, w, h, headers: list, rows: list, col_widths: list = None, header_bg="#1148E8", header_color="#FFFFFF", font_size=11):
        """Add a native editable table."""
        left, top, width, height = self.map_x(x), self.map_y(y), self.map_w(w), self.map_h(h)
        num_rows = len(rows) + 1
        num_cols = len(headers)
        table_shape = self.slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table
        
        if col_widths and len(col_widths) == num_cols:
            for idx, cw in enumerate(col_widths):
                table.columns[idx].width = self.map_w(cw)
                
        for c_idx, head in enumerate(headers):
            cell = table.cell(0, c_idx)
            cell.text = str(head)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(header_bg)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.bold = True
                p.font.color.rgb = hex_to_rgb(header_color)
                p.alignment = PP_ALIGN.CENTER
                
        for r_idx, row in enumerate(rows):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = hex_to_rgb("#333333")
                    p.alignment = PP_ALIGN.CENTER
        return table_shape

    def save(self, out_path: str):
        self.prs.save(out_path)
        print(f"Saved reconstructed presentation to: {out_path}")

if __name__ == "__main__":
    print("PptBuilder helper initialized successfully.")
