#!/usr/bin/env python3
"""Build native editable PPTX from legacy or image2ppt scene v2 JSON."""

from __future__ import annotations

import argparse
import json
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def rgb(value: str) -> RGBColor:
    s = value.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def align(value: str):
    return {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(str(value).lower(), PP_ALIGN.LEFT)


def valign(value: str):
    return {"top": MSO_VERTICAL_ANCHOR.TOP, "middle": MSO_VERTICAL_ANCHOR.MIDDLE, "bottom": MSO_VERTICAL_ANCHOR.BOTTOM}.get(str(value).lower(), MSO_VERTICAL_ANCHOR.MIDDLE)


class Geometry:
    def __init__(self, scene: dict):
        canvas = scene.get("canvas", {})
        self.units = scene.get("units", "in")
        self.source_w = float(canvas.get("width", 1))
        self.source_h = float(canvas.get("height", 1))
        self.slide_w = float(canvas.get("slide_width_in", scene.get("slide_width", 13.333333)))
        self.slide_h = float(canvas.get("slide_height_in", scene.get("slide_height", 7.5)))

    def x(self, value): return float(value) / self.source_w * self.slide_w if self.units == "px" else float(value)
    def y(self, value): return float(value) / self.source_h * self.slide_h if self.units == "px" else float(value)
    def w(self, value): return self.x(value)
    def h(self, value): return self.y(value)


def merged_element(scene: dict, element: dict) -> dict:
    out = deepcopy(scene.get("defaults", {}))
    ref = element.get("style_ref")
    if ref:
        out.update(deepcopy(scene.get("styles", {}).get(ref, {})))
    out.update(deepcopy(element))
    return out


def add_text(slide, e, g: Geometry):
    box = slide.shapes.add_textbox(Inches(g.x(e["x"])), Inches(g.y(e["y"])), Inches(g.w(e["w"])), Inches(g.h(e["h"])))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = bool(e.get("word_wrap", True)); tf.vertical_anchor = valign(e.get("valign", "middle"))
    margin = float(e.get("margin", 0.02))
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.text = str(e.get("text", "")); p.alignment = align(e.get("align", "left"))
    p.font.name = e.get("font", "Microsoft YaHei"); p.font.size = Pt(e.get("size", 12)); p.font.bold = bool(e.get("bold", False)); p.font.color.rgb = rgb(e.get("color", "#0E2240"))
    return box


def add_box(slide, e, g: Geometry):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if e.get("rounded", True) else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(g.x(e["x"])), Inches(g.y(e["y"])), Inches(g.w(e["w"])), Inches(g.h(e["h"])))
    if e.get("fill"):
        shape.fill.solid(); shape.fill.fore_color.rgb = rgb(e["fill"])
    else:
        shape.fill.background()
    if e.get("line"):
        shape.line.color.rgb = rgb(e["line"]); shape.line.width = Pt(e.get("line_width", 1.0))
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, e, g: Geometry):
    shape = slide.shapes.add_connector(1, Inches(g.x(e["x1"])), Inches(g.y(e["y1"])), Inches(g.x(e["x2"])), Inches(g.y(e["y2"])))
    shape.line.color.rgb = rgb(e.get("color", "#0E2240")); shape.line.width = Pt(e.get("width", 1.2))
    return shape


def add_circle(slide, e, g: Geometry):
    r = float(e["r"])
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(g.x(e["cx"] - r)), Inches(g.y(e["cy"] - r)), Inches(g.w(2 * r)), Inches(g.h(2 * r)))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(e.get("fill", "#1148E8"))
    if e.get("line"):
        shape.line.color.rgb = rgb(e["line"]); shape.line.width = Pt(e.get("line_width", 0.5))
    else:
        shape.line.fill.background()
    return shape


def add_table(slide, e, g: Geometry):
    rows, cols = int(e["rows"]), int(e["cols"])
    data = e.get("data", [[""] * cols for _ in range(rows)])
    shape = slide.shapes.add_table(rows, cols, Inches(g.x(e["x"])), Inches(g.y(e["y"])), Inches(g.w(e["w"])), Inches(g.h(e["h"])))
    for row in range(rows):
        for col in range(cols):
            cell = shape.table.cell(row, col)
            cell.text = str(data[row][col]) if row < len(data) and col < len(data[row]) else ""
            cell.fill.solid(); cell.fill.fore_color.rgb = rgb(e.get("header_fill", "#F4F7FB") if row == 0 else e.get("fill", "#FFFFFF"))
            p = cell.text_frame.paragraphs[0]; p.alignment = align(e.get("align", "center")); p.font.name = e.get("font", "Microsoft YaHei"); p.font.size = Pt(e.get("size", 10)); p.font.bold = bool(e.get("header_bold", True) if row == 0 else e.get("bold", False)); p.font.color.rgb = rgb(e.get("color", "#0E2240"))
    return shape


def add_picture(slide, e, g: Geometry, base: Path):
    asset = Path(e["asset"])
    if not asset.is_absolute():
        asset = (base / asset).resolve()
    return slide.shapes.add_picture(str(asset), Inches(g.x(e["x"])), Inches(g.y(e["y"])), Inches(g.w(e["w"])), Inches(g.h(e["h"])))


def render_element(slide, e, g, base):
    funcs = {"text": add_text, "box": add_box, "line": add_line, "circle": add_circle, "table": add_table}
    kind = e["type"].lower()
    if kind == "picture":
        return add_picture(slide, e, g, base)
    if kind not in funcs:
        raise ValueError(f"unsupported element type: {kind}")
    return funcs[kind](slide, e, g)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("spec")
    ap.add_argument("output")
    args = ap.parse_args()
    spec_path = Path(args.spec).resolve()
    scene = json.loads(spec_path.read_text(encoding="utf-8"))
    geometry = Geometry(scene)
    prs = Presentation(); prs.slide_width = Inches(geometry.slide_w); prs.slide_height = Inches(geometry.slide_h)
    slides = scene.get("slides") or [{"background": scene.get("background", "#FFFFFF"), "elements": scene.get("elements", [])}]
    for slide_spec in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(slide_spec.get("background", "#FFFFFF"))
        elements = sorted(slide_spec.get("elements", []), key=lambda item: item.get("z", 0))
        for raw in elements:
            render_element(slide, merged_element(scene, raw), geometry, spec_path.parent)
    output = Path(args.output); output.parent.mkdir(parents=True, exist_ok=True); prs.save(output)
    print(output.resolve())
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

